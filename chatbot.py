import streamlit as st
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
import os
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_google_genai import GoogleGenerativeAIEmbeddings
import google.generativeai as genai
from langchain.vectorstores import FAISS
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain.chains.question_answering import load_qa_chain
from langchain.prompts import PromptTemplate
from dotenv import load_dotenv

load_dotenv()
genai.configure(api_key=os.getenv("GOOGLE_API_KEY"))

# Admin password setup
ADMIN_PASSWORD = "adminpass"

# Function for admin login
def admin_login():
    st.subheader("Admin Login")
    with st.form("login_form"):
        password = st.text_input("Enter admin password:", type="password")
        login_button = st.form_submit_button("Login")
    
        if login_button:
            if password == ADMIN_PASSWORD:
                st.session_state["admin_logged_in"] = True
                st.success("Admin logged in successfully!")
            else:
                st.error("Incorrect password!")

# Admin interface for uploading, processing, deleting, and asking questions
def show_admin_interface():
    st.title("Admin")
    
    st.sidebar.title("Menu:")
    pdf_docs = st.file_uploader("Upload your Files (.pdf, .docx, .ppt, .txt)", accept_multiple_files=True)

    # Ensure file size limit: 2MB per file, 4MB total
    if pdf_docs:
        total_size = sum([file.size for file in pdf_docs])
        if any([file.size > 2 * 1024 * 1024 for file in pdf_docs]):  # 2MB per file
            st.error("Each file must be less than 2MB.")
            return
        if total_size > 4 * 1024 * 1024:  # 4MB total
            st.error("Total file size exceeds 4MB. Please upload smaller files.")
            return

    # Process and store files
    if st.button("Submit & Process"):
        if pdf_docs:
            with st.spinner("Processing..."):
                raw_text = get_text_from_files(pdf_docs)
                text_chunks = get_text_chunks(raw_text)
                get_vector_store(text_chunks, append=True)
                st.session_state["uploaded_files"] = [doc.name for doc in pdf_docs]
                st.success("Documents uploaded and processed!")

    # Display uploaded documents
    with st.container():
        st.subheader("Uploaded Documents:")
        if "uploaded_files" in st.session_state:
            for doc in st.session_state["uploaded_files"]:
                st.write(doc)
        else:
            st.write("No documents uploaded yet.")
    
    # Delete option with space and refresh page after deletion
    st.write("")
    with st.container():
        if "uploaded_files" in st.session_state and st.session_state["uploaded_files"]:
            st.subheader("Delete Uploaded Documents:")
            file_to_delete = st.selectbox("Select a file to delete", st.session_state["uploaded_files"])
            if st.button("Delete File"):
                if file_to_delete in st.session_state["uploaded_files"]:
                    st.session_state["uploaded_files"].remove(file_to_delete)
                    st.success(f"File '{file_to_delete}' deleted successfully!")
                    if st.button("Refresh Page"):
                        # Use query params to trigger a refresh
                        st.experimental_set_query_params()

    # Ask a question with space
    st.write("")
    with st.container():
        st.subheader("Ask a Question from the Uploaded Documents:")
        user_question = st.text_input("Enter your question:")
        if user_question:
            process_user_question(user_question)

# Normal user interface for viewing documents and asking questions
def show_user_interface():
    # st.title("User Interface")
    
    # Display uploaded documents
    with st.container():
        st.subheader("Uploaded Documents:")
        uploaded_files = st.session_state.get("uploaded_files", [])
        if uploaded_files:
            for file in uploaded_files:
                st.write(file)
        else:
            st.write("No documents uploaded yet.")

    # Ask a question
    st.write("")
    with st.container():
        st.subheader("AI Assistant🤖")
        st.subheader("Ask a Question from the Uploaded Documents:")
        user_question = st.text_input("Please enter your question:")
        if user_question:
            process_user_question(user_question)

# Function to extract text from various document types
def get_text_from_files(files):
    text = ""
    for file in files:
        if file.type == "application/pdf":
            pdf_reader = PdfReader(file)
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
        elif file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            doc = Document(file)
            for para in doc.paragraphs:
                text += para.text + "\n"
        elif file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            ppt = Presentation(file)
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        elif file.type == "text/plain":
            text += str(file.read(), "utf-8") + "\n"
    return text

# Function to split text into chunks for processing
def get_text_chunks(text):
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=10000, chunk_overlap=1000)
    return text_splitter.split_text(text)

# Store text chunks in a vector store (FAISS)
def get_vector_store(text_chunks, append=False):
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    
    if append:
        try:
            vector_store = FAISS.load_local("faiss_index", embeddings, allow_dangerous_deserialization=True)
            vector_store.add_texts(text_chunks)
        except Exception as e:
            st.error(f"Error loading FAISS index: {e}")
            vector_store = FAISS.from_texts(text_chunks, embedding=embeddings)
    else:
        vector_store = FAISS.from_texts(text_chunks, embedding=embeddings)

    vector_store.save_local("faiss_index")

# Load the conversational chain
def get_conversational_chain():
    prompt_template = """
    Answer the question as detailed as possible from the provided context. If the answer is not in the provided context, just say, "Answer is not available in the context."
    
    Context:\n {context}?\n
    Question: \n{question}\n

    Answer:
    """
    model = ChatGoogleGenerativeAI(model="gemini-pro", temperature=0.3)
    prompt = PromptTemplate(template=prompt_template, input_variables=["context", "question"])
    return load_qa_chain(model, chain_type="stuff", prompt=prompt)

# Process user questions
def process_user_question(user_question):
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    
    try:
        new_db = FAISS.load_local("faiss_index", embeddings, allow_dangerous_deserialization=True)
        docs = new_db.similarity_search(user_question)
        chain = get_conversational_chain()
        response = chain({"input_documents": docs, "question": user_question}, return_only_outputs=True)
        st.write("Reply: ", response["output_text"])
    except Exception as e:
        st.error("Error: Could not process the question. Ensure documents are uploaded and indexed.")
        st.error(f"Details: {e}")

# Main function
def main():
    st.set_page_config("Document Chat Bot", page_icon="📂")

    st.title("Document Chat Bot")
    st.write("The Document Search Bot assists in managing and searching documents efficiently.")

    if "admin_logged_in" not in st.session_state:
        st.session_state["admin_logged_in"] = False
    if "uploaded_files" not in st.session_state:
        st.session_state["uploaded_files"] = []

    user_type = st.sidebar.radio("Select User Type", ("Admin", "User"))

    if user_type == "Admin":
        if not st.session_state["admin_logged_in"]:
            admin_login()
        else:
            show_admin_interface()
    elif user_type == "User":
        st.session_state["admin_logged_in"] = False
        show_user_interface()

if __name__ == "__main__":
    main()
